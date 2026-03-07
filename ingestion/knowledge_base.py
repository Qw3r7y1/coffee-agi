import io
import json
import os
import uuid
from typing import Optional

from loguru import logger
from rank_bm25 import BM25Okapi


class KnowledgeBase:
    """
    Persistent BM25 knowledge base.
    Documents are stored as JSON on disk and re-indexed on startup.
    """

    INDEX_FILE = "data/chroma/kb_index.json"

    def __init__(self):
        os.makedirs("data/chroma", exist_ok=True)
        self._docs: list[dict] = []   # {id, text, topic, source, difficulty}
        self._bm25: Optional[BM25Okapi] = None
        self._load()

    # ── Public API ────────────────────────────────────────────────────────────

    def search(
        self,
        query: str,
        n_results: int = 4,
        topic_filter: Optional[str] = None,
    ) -> list[dict]:
        if not self._docs or self._bm25 is None:
            return []

        pool = self._docs
        pool_indices = list(range(len(pool)))

        if topic_filter:
            pool_indices = [i for i, d in enumerate(pool) if d["topic"] == topic_filter]
            if not pool_indices:
                return []

        tokens = self._tokenize(query)

        if topic_filter:
            # Re-build a mini BM25 on the filtered subset
            subset_texts = [self._tokenize(self._docs[i]["text"]) for i in pool_indices]
            bm25 = BM25Okapi(subset_texts)
            scores = bm25.get_scores(tokens)
            ranked = sorted(range(len(pool_indices)), key=lambda x: scores[x], reverse=True)
            top_indices = [pool_indices[ranked[i]] for i in range(min(n_results, len(ranked)))]
        else:
            scores = self._bm25.get_scores(tokens)
            top_k = min(n_results, len(self._docs))
            top_indices = sorted(range(len(scores)), key=lambda x: scores[x], reverse=True)[:top_k]

        return [
            {
                "text": self._docs[i]["text"],
                "topic": self._docs[i]["topic"],
                "source": self._docs[i]["source"],
                "difficulty": self._docs[i]["difficulty"],
            }
            for i in top_indices
            if scores[top_indices.index(i) if not topic_filter else pool_indices.index(i)] > 0
        ]

    async def ingest_file(
        self,
        filename: str,
        content: bytes,
        topic_tag: str,
        difficulty_level: str,
    ) -> dict:
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

        if ext == "pdf":
            chunks = self._parse_pdf(content)
        elif ext in ("pptx", "ppt"):
            chunks = self._parse_pptx(content)
        elif ext in ("txt", "md"):
            chunks = self._parse_text(content.decode("utf-8", errors="ignore"))
        else:
            return {"error": f"Unsupported format: .{ext}. Use PDF, PPTX, TXT, or MD."}

        if not chunks:
            return {"error": "No text could be extracted from the file."}

        for chunk in chunks:
            self._docs.append(
                {
                    "id": str(uuid.uuid4()),
                    "text": chunk,
                    "topic": topic_tag,
                    "source": filename,
                    "difficulty": difficulty_level,
                }
            )

        self._rebuild_index()
        self._save()

        logger.info(f"Ingested {len(chunks)} chunks from '{filename}' [{topic_tag} / {difficulty_level}]")
        return {
            "status": "ok",
            "filename": filename,
            "chunks": len(chunks),
            "topic": topic_tag,
            "difficulty": difficulty_level,
        }

    def get_topics(self) -> dict:
        topic_counts: dict[str, int] = {}
        for d in self._docs:
            topic_counts[d["topic"]] = topic_counts.get(d["topic"], 0) + 1
        return {"topics": topic_counts, "total_chunks": len(self._docs)}

    # ── Persistence ───────────────────────────────────────────────────────────

    def _save(self):
        try:
            with open(self.INDEX_FILE, "w", encoding="utf-8") as f:
                json.dump(self._docs, f, ensure_ascii=False)
        except Exception as e:
            logger.warning(f"Failed to save KB index: {e}")

    def _load(self):
        if os.path.exists(self.INDEX_FILE):
            try:
                with open(self.INDEX_FILE, "r", encoding="utf-8") as f:
                    self._docs = json.load(f)
                self._rebuild_index()
                logger.info(f"Knowledge base loaded: {len(self._docs)} chunks")
            except Exception as e:
                logger.warning(f"Failed to load KB index: {e}")
                self._docs = []

    def _rebuild_index(self):
        if self._docs:
            tokenized = [self._tokenize(d["text"]) for d in self._docs]
            self._bm25 = BM25Okapi(tokenized)

    # ── Parsers ───────────────────────────────────────────────────────────────

    def _parse_pdf(self, content: bytes) -> list[str]:
        # 1. Try normal text extraction first
        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(content))
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
            if text.strip():
                return self._chunk_text(text)
        except Exception as e:
            logger.warning(f"PDF text extraction failed, trying OCR: {e}")

        # 2. Fall back to Claude vision OCR for image-based PDFs
        logger.info("PDF has no extractable text — using Claude vision OCR.")
        return self._ocr_pdf_with_claude(content)

    def _ocr_pdf_with_claude(self, content: bytes) -> list[str]:
        try:
            import base64
            import anthropic
            import fitz  # PyMuPDF

            doc = fitz.open(stream=content, filetype="pdf")
            client = anthropic.Anthropic()
            all_text = []

            for page_num, page in enumerate(doc):
                # Render page to PNG at 150 DPI
                mat = fitz.Matrix(150 / 72, 150 / 72)
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes("png")
                img_b64 = base64.standard_b64encode(img_bytes).decode()

                resp = client.messages.create(
                    model="claude-haiku-4-5-20251001",
                    max_tokens=2048,
                    messages=[{
                        "role": "user",
                        "content": [
                            {
                                "type": "image",
                                "source": {"type": "base64", "media_type": "image/png", "data": img_b64},
                            },
                            {
                                "type": "text",
                                "text": "Extract all text from this document page. Output only the raw text, no commentary.",
                            },
                        ],
                    }],
                )
                page_text = resp.content[0].text.strip()
                if page_text:
                    all_text.append(page_text)
                logger.info(f"OCR page {page_num + 1}/{len(doc)}: {len(page_text)} chars")

            doc.close()
            return self._chunk_text("\n\n".join(all_text))
        except Exception as e:
            logger.error(f"Claude OCR failed: {e}")
            return []

    def _parse_pptx(self, content: bytes) -> list[str]:
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            parts = []
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    parts.append(" ".join(slide_text))
            return self._chunk_text("\n".join(parts))
        except Exception as e:
            logger.error(f"PPTX parse error: {e}")
            return []

    def _parse_text(self, text: str) -> list[str]:
        return self._chunk_text(text)

    # ── Helpers ───────────────────────────────────────────────────────────────

    @staticmethod
    def _tokenize(text: str) -> list[str]:
        return text.lower().split()

    @staticmethod
    def _chunk_text(text: str, chunk_size: int = 400, overlap: int = 40) -> list[str]:
        words = text.split()
        if not words:
            return []
        step = max(chunk_size - overlap, 1)
        chunks = []
        for i in range(0, len(words), step):
            chunk = " ".join(words[i : i + chunk_size])
            if chunk.strip():
                chunks.append(chunk)
        return chunks
